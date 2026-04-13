import docx
from PyPDF2 import PdfReader
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import io
import csv
import pandas as pd
from pptx import Presentation

# Set Tesseract path (Windows)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


# -----------------------------
# READ PDF
# -----------------------------
def read_pdf(file):
    """
    Extract text from PDF.
    If no text found (scanned PDF), fallback to OCR.
    """
    text = ""
    try:
        reader = PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print("PDF text extraction error:", e)

    # If PDF has no text -> run OCR
    if not text.strip():
        file.seek(0)
        text = read_pdf_with_ocr(file)
    return text


# -----------------------------
# OCR FOR SCANNED PDF
# -----------------------------
def read_pdf_with_ocr(file):
    """
    Convert PDF pages to images and apply OCR
    """
    text = ""
    try:
        images = convert_from_bytes(file.read())
        for img in images:
            text += pytesseract.image_to_string(img) + "\n"
    except Exception as e:
        print("PDF OCR error:", e)
    return text


# -----------------------------
# READ DOCX, PPTX, CSV, XLS/XLSX
# -----------------------------

def read_docx(file):
    text = ""
    try:
        # Get filename safely
        filename = ""
        if hasattr(file, 'filename') and file.filename:
            filename = file.filename.lower()
        elif hasattr(file, 'name') and file.name:
            filename = file.name.lower()

        if filename.endswith(".docx"):
            doc = docx.Document(file)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        elif filename.endswith((".xls", ".xlsx")):
            df = pd.read_excel(file)
            text += df.to_string(index=False)

        elif filename.endswith(".csv"):
            import io
            if hasattr(file, 'read'):
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode('utf-8',
                        errors='ignore')
                reader = csv.reader(
                    io.StringIO(content))
                for row in reader:
                    text += " | ".join(row) + "\n"

    except Exception as e:
        print(f"Document read error: {e}")
    return text


# -----------------------------
# READ IMAGE (OCR)
# -----------------------------
def read_image(file):
    """
    Extract text from images using OCR
    """
    text = ""
    try:
        img = Image.open(file)
        if img.mode != "RGB":
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img)
    except Exception as e:
        print("Image OCR error:", e)
    return text